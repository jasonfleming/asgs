import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE


# Get command line argument
fname = sys.argv[1]

# Read run.properties and make a property dictionary
runProp = dict()
f = open('run.properties','r')
for line in f:
    fields = line.split(':',1)
    runProp[fields[0].strip()] = fields[1].strip() 
f.close()

# Convert advisoryTime to python datetime object
advisory_dt = datetime.strptime(runProp['time.forecast.valid.cdt'],'%Y%m%d%H%M%S')
advisory_dt_long = datetime.strftime(advisory_dt,'%b-%d-%Y %H:%M')

scenario = runProp['asgs.enstorm']
if [[ scenario == 'nhcConsensus' ]]:
    scenario = 'NHC Track'

prs = Presentation('LSU_template.pptx')

slide_layout = prs.slide_layouts[5]
slide_layout_hydro = prs.slide_layouts[6]

numSlides = 1

# Create a title slide
title_slide_layout = prs.slide_layouts[0]
#for shape in title_slide_layout.placeholders:
#    print('%d %s' % (shape.placeholder_format.idx, shape.name))
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
#title.text = "STORM " + storm + " (" + enstorm + ")"
title.text = runProp['storm class'] + ' ' + runProp['storm name'] + ', ' + scenario + ' Scenario'
subtitle.text = "Advisory " + runProp['advisory'] + " Issued on " + advisory_dt_long + " CDT"
statement = 'For Official Use Only. Not For Release. \rModel results were produced by the ADCIRC Surge Guidance System (ASGS) and are based on the National Hurricane Center (NHC) forecast track.'
fouo = slide.placeholders[10]
fouo.text = statement
numSlides = numSlides + 1

# Set slide layout
left = Inches(1.94)
top = Inches(1.14)

img_path = fname
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'NHC Advisory ' + runProp['advisory'] + ' ' + scenario + ' Scenario'
subtitle.text = "Simulated peak water levels (ft, NAVD88)"
pic = slide.shapes.add_picture(img_path,left,top)
fouo = slide.placeholders[13]
fouo.text = statement
snum = slide.placeholders[14]
snum.text = str(numSlides)
numSlides = numSlides + 1
#for shape in slide.placeholders:
#    print('%d %s' % (shape.placeholder_format.idx, shape.name))

#left = Inches(0.42)
#top = Inches(1.15)
left = Inches(0.75)
top = Inches(0.81)
iwidth = Inches(11.84)
iheight = Inches(5.69)

fnames = ['WSE_17StCanal_USACE85625.png','WSE_IHNC01_USACE76065.png','WSE_IHNC02_USACE76030.png',
        'WSE_LPV144_USACE76010.png','WSE_LPV149_USACE85760.png','WSE_NOV13_USACE01440.png',
        'WSE_NOV14_USACE01440.png','WSE_WBV09a_USACE82770.png','WSE_WBV09b_USACE82762.png',
        'WSE_WBV162_USACE82742.png','WSE_WBV7274_USACE82715.png','WSE_WBV90_USACE76265.png',
        'WSE_LakefrontAirport_USACE85670.png','WSE_Mandeville_USACE85575.png',
        'WSE_Rigolets_USACE85700.png','WSE_Lafitte_USACE82875.png']

# Station names correspond to the order of fnames
#staName = ['17th St. Outfall Canal','Seabrook Complex (IHNC-01)','IHNC Surge Barrier (IHNC-02)',
#        'Bayou Dupre Sector Gate (LPV-144)','Caernarvon Canal Sector Gate (LPV-149)',
#        'Empire Floodgate (NOV-13)','Empire Lock (NOV-14)','Oakville Sluice Gate (WBV-09a)',
#        'Hero Canal stop-log gage (WBV-09b)','Bayou Segnetee closure (WBV-16.2)',
#        'Western Tie-In features (WBV-74-72)','West Closure Complex (WBV-90)',
#        'Lakefront Airport','Mandeville','Rigolets','Lafitte']
staName = ['17th St. Outfall Canal','Seabrook Complex (IHNC01)','IHNC Surge Barrier (IHNC02)',
        'Bayou Dupre Sector Gate (LPV144)','Caernarvon Canal Sector Gate (LPV149)',
        'Empire Floodgate (NOV13)','Empire Lock (NOV14)','Oakville Sluice Gate (WBV09a)',
        'Hero Canal stop-log gage (WBV09b)','Bayou Segnetee closure (WBV162)',
        'Western Tie-In features (WBV7472)','West Closure Complex (WBV90)',
        'Lakefront Airport','Mandeville','Rigolets','Lafitte']

i = 0
for image in fnames:
    try: 
        slide = prs.slides.add_slide(slide_layout_hydro)
        title = slide.shapes.title
        title.text = staName[i]
        pic = slide.shapes.add_picture(image,left,top,width=iwidth,height=iheight)
        fouo = slide.placeholders[13]
        fouo.text = statement
        snum = slide.placeholders[14]
        snum.text = str(numSlides)
        numSlides = numSlides + 1
        i = i + 1
    except:
        print("ERROR: buildPPT.py: Could not find " + image + ".")

# Loop through slides
#slides = prs.slides
#for slide in slides:
        #print('slide number %s' % str(slides.index(slide)+1))

pptFile = runProp['storm name'] + "_Adv" + runProp['advisory'] + "_" + scenario + "_" + runProp['forecastValidStart'] + ".pptx"
prs.save(pptFile)
pFile = open('pptFile.temp','w')
pFile.write(pptFile)
pFile.close()
